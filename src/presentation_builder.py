from __future__ import annotations
import logging
from pathlib import Path
from typing import Dict, Any, Optional, Union, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.dml.color import RGBColor

from .models import PresentationOutline, SlideOutline, SlideType

logger = logging.getLogger(__name__)


class PresentationBuilder:
    """
    Builds PPTX presentations from outlines. 
    Supports optional template loading and theming, with robust handling of common template filler text.
    """

    LAYOUT_MAP: Dict[SlideType, int] = {
        SlideType.TITLE_SLIDE: 0,
        SlideType.CONTENT_SLIDE: 1,
        SlideType.TWO_COLUMN: 2,
        SlideType.CLOSING_SLIDE: 3,
    }

    FILLER_SNIPPETS = (
        "Click to edit Master text styles",
        "Lorem ipsum",
        "Text here",
    )

    def __init__(self, template_config: Optional[Union[Dict[str, Any], str, Path]] = None):
        self.template_config: Dict[str, Any] = {}
        template_path: Optional[Path] = None

        if isinstance(template_config, dict):
            self.template_config = template_config
            p = template_config.get("path")
            if p:
                template_path = Path(p)
        elif isinstance(template_config, (str, Path)):
            template_path = Path(template_config)

        if template_path and template_path.exists():
            self.prs = Presentation(str(template_path))
        else:
            self.prs = Presentation()

        # Remove all existing slides from template deck while preserving masters/layouts
        self._remove_existing_slides()

        # Theme (optional); fallback defaults
        self.theme = (self.template_config.get("theme") or {})
        self._apply_default_theme()

    # ----------------------------
    # Template housekeeping
    # ----------------------------
    def _remove_existing_slides(self) -> None:
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

    # ----------------------------
    # Theme helpers
    # ----------------------------
    def _apply_default_theme(self) -> None:
        colors = self.theme.get("colors") or {}
        sizes = self.theme.get("sizes") or {}

        colors.setdefault("primary", "#3A0CA3")
        colors.setdefault("accent", "#7209B7")
        colors.setdefault("background", "#FFFFFF")
        colors.setdefault("subtle_bg", "#F2F2F2")
        colors.setdefault("text_dark", "#111111")
        colors.setdefault("text_light", "#FFFFFF")

        sizes.setdefault("heading_large", 36)
        sizes.setdefault("heading_medium", 28)
        sizes.setdefault("heading_small", 18)
        sizes.setdefault("body_text", 16)

        self.theme["colors"] = colors
        self.theme["sizes"] = sizes

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        hex_color = hex_color.lstrip("#")
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )

    # ----------------------------
    # Layout + placeholder helpers
    # ----------------------------
    def _safe_layout(self, slide_type: SlideType):
        desired = self.LAYOUT_MAP.get(slide_type, 1)
        max_idx = len(self.prs.slide_layouts) - 1
        idx = desired if 0 <= desired <= max_idx else min(1, max_idx)
        return self.prs.slide_layouts[idx]

    def _get_placeholder(self, prs_slide, placeholder_type: PP_PLACEHOLDER):
        for shape in prs_slide.placeholders:
            try:
                if shape.placeholder_format.type == placeholder_type:
                    return shape
            except Exception:
                continue
        return None

    def _clear_template_filler(self, prs_slide) -> None:
        """
        Clears any text in shapes that matches known template filler snippets.
        """
        for shape in prs_slide.shapes:
            try:
                if not getattr(shape, "has_text_frame", False):
                    continue
                txt = (shape.text_frame.text or "").strip()
                if not txt:
                    continue
                if any(snippet in txt for snippet in self.FILLER_SNIPPETS):
                    shape.text_frame.clear()
            except Exception:
                # Keep generation robust across odd shapes
                continue

    def _set_textframe(
        self,
        text_frame,
        lines: List[str],
        *,
        font_size: int,
        bold_first: bool = False,
        color: Optional[RGBColor] = None,
        align: Optional[int] = None,
    ) -> None:
        text_frame.clear()
        text_frame.word_wrap = True

        for i, line in enumerate(lines):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(font_size)
            if i == 0 and bold_first:
                p.font.bold = True
            if color is not None:
                p.font.color.rgb = color
            if align is not None:
                p.alignment = align

    def _add_background_rect(self, prs_slide, hex_color: str) -> None:
        bg = prs_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.prs.slide_width,
            self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb(hex_color)
        bg.line.fill.background()

        prs_slide.shapes._spTree.remove(bg._element)
        prs_slide.shapes._spTree.insert(2, bg._element)

    # ----------------------------
    # Slide builders
    # ----------------------------
    def _add_title_slide(self, slide: SlideOutline) -> None:
        prs_slide = self.prs.slides.add_slide(self._safe_layout(slide.slide_type))

        # Clear template filler BEFORE adding our content
        self._clear_template_filler(prs_slide)

        # Background (optional)
        self._add_background_rect(prs_slide, self.theme["colors"]["primary"])

        title_ph = self._get_placeholder(prs_slide, PP_PLACEHOLDER.TITLE)
        if title_ph is not None:
            self._set_textframe(
                title_ph.text_frame,
                [slide.title],
                font_size=self.theme["sizes"]["heading_large"],
                bold_first=True,
                color=self._hex_to_rgb(self.theme["colors"]["text_light"]),
                align=PP_ALIGN.CENTER,
            )
        else:
            box = prs_slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(1.5))
            self._set_textframe(
                box.text_frame,
                [slide.title],
                font_size=self.theme["sizes"]["heading_large"],
                bold_first=True,
                color=self._hex_to_rgb(self.theme["colors"]["text_light"]),
                align=PP_ALIGN.CENTER,
            )

        if slide.subtitle:
            sub_ph = self._get_placeholder(prs_slide, PP_PLACEHOLDER.SUBTITLE)
            if sub_ph is not None:
                self._set_textframe(
                    sub_ph.text_frame,
                    [slide.subtitle],
                    font_size=self.theme["sizes"]["heading_small"],
                    color=self._hex_to_rgb(self.theme["colors"]["text_light"]),
                    align=PP_ALIGN.CENTER,
                )
            else:
                box = prs_slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(9), Inches(1.0))
                self._set_textframe(
                    box.text_frame,
                    [slide.subtitle],
                    font_size=self.theme["sizes"]["heading_small"],
                    color=self._hex_to_rgb(self.theme["colors"]["text_light"]),
                    align=PP_ALIGN.CENTER,
                )

        # Clear template filler AFTER as well (in case subtitle placeholder is included in the layout and contains filler text)
        self._clear_template_filler(prs_slide)

    def _add_content_slide(self, slide: SlideOutline) -> None:
        prs_slide = self.prs.slides.add_slide(self._safe_layout(slide.slide_type))

        # Clear filler immediately
        self._clear_template_filler(prs_slide)

        # Title
        title_ph = self._get_placeholder(prs_slide, PP_PLACEHOLDER.TITLE)
        if title_ph is not None:
            self._set_textframe(
                title_ph.text_frame,
                [slide.title],
                font_size=self.theme["sizes"]["heading_medium"],
                bold_first=True,
                color=self._hex_to_rgb(self.theme["colors"]["primary"]),
            )
        else:
            box = prs_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            self._set_textframe(
                box.text_frame,
                [slide.title],
                font_size=self.theme["sizes"]["heading_medium"],
                bold_first=True,
                color=self._hex_to_rgb(self.theme["colors"]["primary"]),
            )

        bullets = [b for b in (slide.bullet_points or []) if isinstance(b, str) and b.strip()]
        if bullets:
            body_ph = self._get_placeholder(prs_slide, PP_PLACEHOLDER.BODY)
            if body_ph is not None:
                self._set_textframe(
                    body_ph.text_frame,
                    bullets,
                    font_size=self.theme["sizes"]["body_text"],
                    color=self._hex_to_rgb(self.theme["colors"]["text_dark"]),
                )
            else:
                box = prs_slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
                self._set_textframe(
                    box.text_frame,
                    bullets,
                    font_size=self.theme["sizes"]["body_text"],
                    color=self._hex_to_rgb(self.theme["colors"]["text_dark"]),
                )

        # Clear filler again (handles extra textboxes included in the layout)
        self._clear_template_filler(prs_slide)

    def _add_closing_slide(self, slide: SlideOutline) -> None:
        prs_slide = self.prs.slides.add_slide(self._safe_layout(slide.slide_type))

        self._clear_template_filler(prs_slide)
        self._add_background_rect(prs_slide, self.theme["colors"]["accent"])

        title_ph = self._get_placeholder(prs_slide, PP_PLACEHOLDER.TITLE)
        if title_ph is not None:
            self._set_textframe(
                title_ph.text_frame,
                [slide.title],
                font_size=self.theme["sizes"]["heading_large"],
                bold_first=True,
                color=self._hex_to_rgb(self.theme["colors"]["text_light"]),
                align=PP_ALIGN.CENTER,
            )
        else:
            box = prs_slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(1.5))
            self._set_textframe(
                box.text_frame,
                [slide.title],
                font_size=self.theme["sizes"]["heading_large"],
                bold_first=True,
                color=self._hex_to_rgb(self.theme["colors"]["text_light"]),
                align=PP_ALIGN.CENTER,
            )

        if slide.subtitle:
            box = prs_slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(9), Inches(1.0))
            self._set_textframe(
                box.text_frame,
                [slide.subtitle],
                font_size=self.theme["sizes"]["heading_small"],
                color=self._hex_to_rgb(self.theme["colors"]["text_light"]),
                align=PP_ALIGN.CENTER,
            )

        bullets = [b for b in (slide.bullet_points or []) if isinstance(b, str) and b.strip()]
        if bullets:
            body_ph = self._get_placeholder(prs_slide, PP_PLACEHOLDER.BODY)
            if body_ph is not None:
                self._set_textframe(
                    body_ph.text_frame,
                    bullets,
                    font_size=self.theme["sizes"]["body_text"],
                    color=self._hex_to_rgb(self.theme["colors"]["text_light"]),
                    align=PP_ALIGN.CENTER,
                )

        self._clear_template_filler(prs_slide)

    # ----------------------------
    # Public API
    # ----------------------------
    def build_from_outline(self, outline: PresentationOutline) -> Presentation:
        for slide in outline.slides:
            if slide.slide_type == SlideType.TITLE_SLIDE:
                self._add_title_slide(slide)
            elif slide.slide_type == SlideType.CLOSING_SLIDE:
                self._add_closing_slide(slide)
            else:
                self._add_content_slide(slide)
            logger.info(f"Added slide {slide.slide_number}: {slide.title}")
        return self.prs

    def save(self, output_path: Path) -> str:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        logger.info(f"Presentation saved to {output_path}")
        return str(output_path)