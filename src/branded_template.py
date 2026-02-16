from __future__ import annotations

import logging
from pathlib import Path
from typing import List, Optional, Dict, Any, Iterable, Tuple

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from .models import SlideType

logger = logging.getLogger(__name__)


# ------------------------------------------------------------
# Template handler
# ------------------------------------------------------------
class BrandedTemplateHandler:
    def __init__(self, template_path: Path):
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        self.template_path = template_path

    def create_presentation_from_template(self) -> Presentation:
        return Presentation(str(self.template_path))

    # Avoid layouts that create grid/table look
    AVOID = ("onepager", "grid", "table")

    PREFERRED = {
        SlideType.TITLE_SLIDE: ("cover", "title"),
        SlideType.CONTENT_SLIDE: ("agenda", "content", "key_message"),
        SlideType.TWO_COLUMN: ("two", "comparison"),
        SlideType.CLOSING_SLIDE: ("closing", "summary", "key_message"),
    }

    def pick_layout_index(self, prs: Presentation, slide_type: SlideType) -> int:
        preferred = self.PREFERRED.get(slide_type, ("content",))

        for i, layout in enumerate(prs.slide_layouts):
            name = (layout.name or "").lower()
            if any(bad in name for bad in self.AVOID):
                continue
            if any(good in name for good in preferred):
                return i

        return 0  # safe fallback


# ------------------------------------------------------------
# Content injector
# ------------------------------------------------------------
class TemplateContentInjector:
    FILLER_MARKERS = (
        "click to add title",
        "click to add text",
        "first level",
        "second level",
        "third level",
        "lorem ipsum",
    )

    def __init__(self, handler: BrandedTemplateHandler):
        self.handler = handler

    # ---------- helpers ----------
    def _iter_shapes(self, shapes):
        for s in shapes:
            yield s
            if hasattr(s, "shapes"):
                yield from self._iter_shapes(s.shapes)

    def _clear_filler(self, slide):
        for s in self._iter_shapes(slide.shapes):
            if not getattr(s, "has_text_frame", False):
                continue
            txt = (s.text_frame.text or "").lower()
            if any(m in txt for m in self.FILLER_MARKERS):
                s.text_frame.clear()

    def _get_placeholders(self, slide, types: Tuple[PP_PLACEHOLDER, ...]):
        found = []
        for s in slide.placeholders:
            try:
                if s.placeholder_format.type in types and s.has_text_frame:
                    found.append(s)
            except Exception:
                pass
        return sorted(found, key=lambda x: x.placeholder_format.idx)

    def _fill(self, tf, lines: List[str]):
        tf.clear()
        if not lines:
            return
        tf.paragraphs[0].text = lines[0]
        for l in lines[1:]:
            p = tf.add_paragraph()
            p.text = l
            p.level = 0

    # ---------- public ----------
    def inject_content(self, prs: Presentation, outlines: List[Any]) -> Presentation:
        while prs.slides:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        for outline in outlines:
            self._add_slide(prs, outline)

        return prs

    def _add_slide(self, prs: Presentation, outline: Any):
        slide_type = outline.slide_type
        layout_idx = self.handler.pick_layout_index(prs, slide_type)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        self._clear_filler(slide)

        # -------- TITLE --------
        title_ph = self._get_placeholders(
            slide, (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
        )
        if title_ph:
            self._fill(title_ph[0].text_frame, [outline.title])

        # -------- BODY / OBJECT --------
        bullets = list(outline.bullet_points or [])

        body_ph = self._get_placeholders(slide, (PP_PLACEHOLDER.BODY,))
        obj_ph = self._get_placeholders(slide, (PP_PLACEHOLDER.OBJECT,))

        if body_ph:
            self._fill(body_ph[0].text_frame, bullets)
            for extra in body_ph[1:]:
                extra.text_frame.clear()
            for o in obj_ph:
                o.text_frame.clear()
        elif obj_ph:
            self._fill(obj_ph[0].text_frame, bullets)
            for extra in obj_ph[1:]:
                extra.text_frame.clear()

        self._clear_filler(slide)